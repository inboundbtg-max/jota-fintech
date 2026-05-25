# 📊 GERADOR DE PPTX PROFISSIONAL
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
from logger import setup_logger

logger = setup_logger(__name__)


class PPTXGenerator:
    """
    Gera apresentação PowerPoint profissional a partir do relatório.

    Padrão senior: esquemas de cores, estrutura clara, conteúdo bem organizado.
    """

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

        # Cores corporativas
        self.COLOR_PRIMARY = RGBColor(0, 51, 102)  # Azul escuro
        self.COLOR_ACCENT = RGBColor(0, 102, 204)  # Azul claro
        self.COLOR_TEXT = RGBColor(64, 64, 64)     # Cinza escuro

        logger.info("✅ Gerador PPTX inicializado")

    def add_title_slide(self, title: str, subtitle: str = ""):
        """Adiciona slide de título"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Background
        background = slide.shapes.add_shape(
            1,  # Rectangle
            0, 0,
            self.prs.slide_width,
            self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLOR_PRIMARY
        background.line.color.rgb = self.COLOR_PRIMARY

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(2.5),
            Inches(9),
            Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(4.2),
                Inches(9),
                Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = subtitle
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = self.COLOR_ACCENT
            subtitle_para.alignment = PP_ALIGN.CENTER

        # Data
        date_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(6.8),
            Inches(9),
            Inches(0.5)
        )
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_para.text = f"Data: {datetime.now().strftime('%d de %B de %Y')}"
        date_para.font.size = Pt(14)
        date_para.font.color.rgb = RGBColor(200, 200, 200)
        date_para.alignment = PP_ALIGN.CENTER

        logger.info("✅ Slide de título adicionado")

    def add_content_slide(self, title: str, content: str):
        """Adiciona slide com conteúdo (título + bullets)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Fundo branco
        background = slide.shapes.add_shape(
            1,
            0, 0,
            self.prs.slide_width,
            self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        background.line.color.rgb = RGBColor(200, 200, 200)

        # Header bar
        header = slide.shapes.add_shape(
            1,
            0, 0,
            self.prs.slide_width,
            Inches(0.8)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.COLOR_PRIMARY
        header.line.color.rgb = self.COLOR_PRIMARY

        # Título
        title_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.15),
            Inches(9),
            Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Conteúdo
        content_box = slide.shapes.add_textbox(
            Inches(0.7),
            Inches(1.2),
            Inches(8.6),
            Inches(5.8)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        # Adiciona conteúdo (com quebras de linha como bullets)
        lines = content.strip().split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = line.strip()
            p.font.size = Pt(16)
            p.font.color.rgb = self.COLOR_TEXT
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)

        logger.info(f"✅ Slide de conteúdo adicionado: {title}")

    def add_summary_slide(self, summary: str):
        """Adiciona slide de resumo executivo"""
        self.add_content_slide("Resumo Executivo", summary)

    def save(self, filename: str = "relatorio_estrategico.pptx") -> str:
        """Salva a apresentação"""
        try:
            self.prs.save(filename)
            logger.info(f"✅ PPTX salvo: {filename}")
            return filename
        except Exception as e:
            logger.error(f"❌ Erro ao salvar PPTX: {str(e)}")
            raise


def generate_pptx_from_report(briefing: str, report: str, filename: str = "relatorio_estrategico.pptx") -> str:
    """
    Função helper para gerar PPTX completo a partir de um relatório.

    Args:
        briefing: Briefing original do cliente
        report: Relatório consolidado dos agentes
        filename: Nome do arquivo PPTX a salvar

    Returns:
        str: Caminho do arquivo gerado
    """
    logger.info("📊 Gerando PPTX...")

    generator = PPTXGenerator()

    # Slide 1: Título
    generator.add_title_slide(
        title="Planejamento Estratégico",
        subtitle="Proposta de Estratégia Digital"
    )

    # Slide 2: Briefing
    generator.add_content_slide(
        title="Briefing do Cliente",
        content=briefing.strip()
    )

    # Slide 3: Relatório
    generator.add_content_slide(
        title="Análise Estratégica",
        content=report.strip()
    )

    # Slide 4: Próximos Passos
    generator.add_content_slide(
        title="Próximos Passos",
        content="""
• Validação da estratégia com stakeholders
• Desenvolvimento do plano detalhado
• Implementação das ações recomendadas
• Monitoramento e otimização contínua
• Relatórios periódicos de performance
        """.strip()
    )

    # Salvar
    filepath = generator.save(filename)
    logger.info(f"📄 PPTX gerado com sucesso: {filepath}")

    return filepath
